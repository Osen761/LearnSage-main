import os
from os import listdir
from os.path import isfile, join, isdir
import PyPDF2
import docx
from pptx import Presentation
from dotenv import load_dotenv
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_qdrant import Qdrant
from langchain_text_splitters import TokenTextSplitter
import streamlit as st

# Load environment variables
load_dotenv()
api_key = os.environ["GOOGLE_API_KEY"] = st.secrets["GOOGLE_API_KEY"]

class Document:
    def __init__(self, page_content, metadata):
        self.page_content = page_content
        self.metadata = metadata

class DocumentIndexer:
    def __init__(self, model_name="models/embedding-001", local_qdrant_path="output", collection_name="MyCollection"):
        self.model_name = model_name
        self.local_qdrant_path = local_qdrant_path
        self.collection_name = collection_name
        self.embedding_model = GoogleGenerativeAIEmbeddings(model=model_name, google_api_key=api_key)
        self.qdrant = None

    def get_files(self, dir):
        file_list = []
        for f in listdir(dir):
            full_path = join(dir, f)
            if isfile(full_path):
                file_list.append(full_path)
            elif isdir(full_path):
                file_list.extend(self.get_files(full_path))
        return file_list

    def get_text_from_word(self, filename):
        doc = docx.Document(filename)
        return '\n'.join(para.text for para in doc.paragraphs)

    def get_text_from_pptx(self, filename):
        prs = Presentation(filename)
        return '\n'.join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

    def extract_text(self, file):
        if file.endswith(".pdf"):
            reader = PyPDF2.PdfReader(file)
            return " ".join(page.extract_text() for page in reader.pages)
        elif file.endswith(".txt"):
            with open(file, 'r') as f:
                return f.read()
        elif file.endswith(".docx"):
            return self.get_text_from_word(file)
        elif file.endswith(".pptx"):
            return self.get_text_from_pptx(file)
        else:
            return None

    def index_documents(self, mypath):
        onlyfiles = self.get_files(mypath)
        text_splitter = TokenTextSplitter(chunk_size=1000, chunk_overlap=200)
        
        for file in onlyfiles:
            print(f"Indexing {file}")
            file_content = self.extract_text(file)
            if not file_content:
                continue

            texts = text_splitter.split_text(file_content)
            metadata = {"path": file}
            documents = [Document(text, metadata) for text in texts]

            if self.qdrant is None:
                self.qdrant = Qdrant.from_documents(
                    documents,
                    self.embedding_model,
                    path=self.local_qdrant_path,
                    collection_name=self.collection_name
                )
            else:
                self.qdrant.add_documents(documents)

        print(onlyfiles)
        print("Finished indexing!")


